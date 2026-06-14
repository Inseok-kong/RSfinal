"""
Auto-generate a 13-slide PowerPoint deck + Q&A document.

Slides:
  1  Title / Background
  2  Research Questions
  3  Study Area
  4  Data
  5  Analysis Workflow
  6  RGB Comparison
  7  NDVI Result
  8  NDBI Result
  9  MNDWI + BSI Result
 10  Change Detection
 11  Land Cover Classification
 12  Quantitative Analysis
 13  Conclusions

Outputs to outputs/presentation/.
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import config


# ----------------------------------------------------------------------------
# Q&A doc — 10 expected professor questions
# ----------------------------------------------------------------------------
QA_ITEMS = [
    ("Q1. 왜 Landsat Collection 2 Level-2 Surface Reflectance 데이터를 사용했는가?",
     "Landsat C2 L2 데이터는 USGS에서 표준 대기보정(Atmospheric Correction)을 완료하여 "
     "제공하는 표면반사도(Surface Reflectance) 데이터입니다. 서로 다른 시기의 영상을 비교할 때, "
     "대기 조건(에어로졸, 수증기 등)에 의한 오차를 최소화하고 지표면 고유의 분광 특성만을 비교하기 "
     "위해서는 대기보정이 필수적입니다. 또한, 1970년대부터 축적된 동일 센서 설계 사상을 계승하고 있어 "
     "20년 이상의 장기 시계열 분석에 가장 적합합니다."),
    ("Q2. 2000, 2010, 2022년 시기 선정 배경과 계절 통제 이유는 무엇인가?",
     "2000년은 월드컵 유치 발표 이전의 '원시 사막 Baseline', 2010년은 월드컵 유치가 확정되어 "
     "루사일 신도시 개발이 시작된 시점, 2022년은 월드컵 개최 및 신도시 완공 시점입니다. "
     "또한, 태양 고도각 변화와 계절적 식생 및 수분 변화로 인한 오차를 통제하기 위해 모든 영상을 "
     "1~4월(카타르의 건기 및 겨울~봄철)로 일치시켜 분광 반사율 비교의 신뢰성을 확보했습니다."),
    ("Q3. Landsat 7 SLC-off 현상에 따른 데이터 결측을 어떻게 해결했는가?",
     "Landsat 7 ETM+ 센서는 2003년 5월 Scan Line Corrector(SLC) 고장으로 인해 영상의 좌우 외곽부에 "
     "빗살무늬 형태의 데이터 공백(결측값)이 발생합니다. 이를 극복하기 위해, 2010년 분석 시 단일 장면이 "
     "아닌 동일 계절 범위(1~4월) 내 구름 10% 이하인 Landsat 7 영상 6장을 수집하여, 픽셀 단위로 "
     "중간값을 취하는 Median Composite 기법을 사용했습니다. 이를 통해 결측 영역을 다른 날짜의 유효 픽셀로 "
     "자연스럽게 보완(Gap-filling)했습니다."),
    ("Q4. 분석에 사용된 4가지 스펙트럴 지수의 분광학적 원리는 무엇인가?",
     "첫째, NDVI는 식생의 Red 흡수와 NIR 강반사 특성을 이용하여 식생 강도를 평가합니다. "
     "둘째, NDBI는 인공구조물이 SWIR1 대역에서 NIR보다 강한 반사율을 보이는 점을 활용해 시가지를 강조합니다. "
     "셋째, MNDWI는 수체가 Green 대역에서 높은 반사율을 보이고 SWIR1 대역에서 거의 흡수(0에 수렴)되는 원리로 "
     "수체 영역을 강조합니다. 넷째, BSI는 Red와 SWIR1 대역의 높은 흙 반사율과 Blue/NIR의 흡수 특성을 조합하여 "
     "나지 및 모래사막을 추출합니다."),
    ("Q5. 사막과 시가지(Urban)를 구분하기 위해 여러 지수를 함께 사용한 이유는 무엇인가?",
     "사막과 건조기 시가지(아스팔트, 콘크리트 등)는 식생이 없기 때문에 NDVI가 둘 다 매우 낮게 나타납니다. "
     "따라서 NDVI만으로는 둘을 분리할 수 없습니다. 반면, 인공 구조물은 SWIR1 대역의 반사가 도드라져 NDBI가 높고, "
     "천연 사막 모래는 철산화물이나 실리카 성분으로 인해 BSI가 도드라지게 높습니다. 이처럼 다중 분광 채널 및 "
     "다양한 지수(NDVI, NDBI, MNDWI, BSI)를 융합함으로써 혼재된 건조 지역 내 물적 특성을 명확히 분리할 수 있었습니다."),
    ("Q6. KMeans 클러스터링의 작동 원리와 10차원 특징 공간 설계의 장점은 무엇인가?",
     "KMeans는 고차원 특징 공간에서 데이터 포인트들 사이의 유클리드 거리를 최소화하여 k개의 군집으로 묶는 "
     "대표적인 비지도 기법입니다. 본 분석에서는 Landsat의 6개 물리적 반사도 밴드(Blue, Green, Red, NIR, "
     "SWIR1, SWIR2)와 이로부터 추출된 4개 지수(NDVI, NDBI, MNDWI, BSI)를 결합하여 총 10차원의 특징 공간을 "
     "설계했습니다. 단순히 원시 반사도 밴드만 쓰는 것보다 타깃 피복(식생, 도시, 수체, 사막)의 물리적 특성을 "
     "극대화한 지수를 포함시킴으로써 분류 경계가 훨씬 명확해지고 노이즈가 억제되는 효과를 얻었습니다."),
    ("Q7. 세 시기(2000, 2010, 2022)의 KMeans 분류 일관성을 확보한 방법은 무엇인가?",
     "각 연도마다 KMeans 모델을 개별 학습(Fit)시키면 동일한 토지 피복이라도 연도별 클러스터 중심값(Centroid)과 "
     "인덱스 번호가 달라져 직접적인 시계열 비교가 불가능해집니다. 이를 해결하기 위해 세 시기 전체 이미지에서 "
     "픽셀 데이터를 추출하여 연합 풀(Union Pool)을 구성하고, 이에 대해 단일 KMeans 모델을 적합시켰습니다. "
     "이후 이 단일 모델로 각 시기의 토지 피복을 예측(Predict)함으로써 시기 간 완벽한 분류 일관성을 보장했습니다."),
    ("Q8. UTM Zone 39N (EPSG:32639) 투영 좌표계를 선정한 이유와 면적 신뢰성은 무엇인가?",
     "WGS84(EPSG:4326)와 같은 경위도 좌표계는 각도 단위를 사용하므로 정량적인 거리나 면적($\text{km}^2$)을 "
     "계산할 때 위도에 따른 왜곡이 매우 심해집니다. 따라서 분석 대상 지역인 카타르 루사일이 속한 횡단 메르카토르 "
     "투영계인 UTM Zone 39N(EPSG:32639)을 사용하여 30 m 격자 크기가 평면상에서 등면적성을 유지하도록 투영했습니다. "
     "이를 통해 격자 셀 면적($30 \text{m} \times 30 \text{m} = 900 \text{m}^2$)에 기반한 정량적 면적 통계의 신뢰성을 완벽히 확보했습니다."),
    ("Q9. 토지 피복 전이 행렬(Transition Matrix)이 가지는 학술적 및 분석적 의의는 무엇인가?",
     "단순히 연도별 토지 피복의 전체 면적 합계만 비교하면 'A가 줄고 B가 늘었다'는 사실만 알 수 있을 뿐, "
     "실제 어느 위치의 어떤 피복이 다른 피복으로 바뀌었는지는 추적할 수 없습니다. 전이 행렬은 공간 격자 매칭을 통해 "
     "2000년의 사막(Desert)이 2022년에 구체적으로 몇 $\text{km}^2$ 만큼 시가지(Urban)로 바뀌고 몇 $\text{km}^2$ 만큼 "
     "녹지(Vegetation)로 전이되었는지 흐름을 정량 추적할 수 있게 하여, 공간적 변화 경로를 완벽히 증명해 줍니다."),
    ("Q10. 본 연구의 한계점과 이를 원격탐사 관점에서 개선하기 위한 향후 대책은 무엇인가?",
     "첫째, 비지도 분류 기법인 KMeans의 특성상 모래와 도심 인프라가 혼재된 '혼합 지표면(Mixed Surface)'의 "
     "경계부에서 오분류가 일부 발생했습니다. 둘째, 검증을 위한 지상 실측 데이터(Ground Truth)나 고해상도(예: 1m 급) "
     "참조 지도가 없어 정밀 오차 행렬(Confusion Matrix)을 산출하지 못했습니다. 향후 Sentinel-2(10m 해상도)나 "
     "고해상도 위성 자료를 융합하고, 감독 분류(Random Forest) 또는 U-Net과 같은 딥러닝 기반 세그멘테이션 기법을 도입하여 "
     "참조 자료와 교차 검증을 수행한다면 분류 정확도를 더욱 크게 개선할 수 있습니다."),
]


def _add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def _add_content_slide(prs, title, bullets,
                       image_path: Optional[Path] = None,
                       image_left=Inches(5.2), image_top=Inches(1.4),
                       image_height=Inches(5.0)):
    layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Bullet text box (left half)
    left = Inches(0.4); top = Inches(1.4)
    width = Inches(4.7); height = Inches(5.5)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT

    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path),
                                 image_left, image_top,
                                 height=image_height)
    return slide


def _add_image_only_slide(prs, title, image_path: Path):
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if image_path and Path(image_path).exists():
        from PIL import Image
        with Image.open(image_path) as img:
            w_px, h_px = img.size
        aspect = w_px / h_px
        
        # If the image is wide (e.g., triptych comparisons), fit to slide width
        if aspect > 1.5:
            width = Inches(12.33)
            height = width / aspect
            # Vertically center within the available body area (1.3 to 6.9 Inches)
            top = Inches(1.3) + (Inches(5.6) - height) / 2
            slide.shapes.add_picture(str(image_path),
                                     Inches(0.5), top,
                                     width=width)
        else:
            slide.shapes.add_picture(str(image_path),
                                     Inches(0.5), Inches(1.3),
                                     height=Inches(5.6))
    return slide


def _df_table(slide, df: pd.DataFrame,
              left=Inches(0.5), top=Inches(1.4),
              width=Inches(9), height=Inches(5)):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(name)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            v = df.iat[i, j]
            cell = table.cell(i + 1, j)
            if isinstance(v, float):
                cell.text = f"{v:.3f}"
            else:
                cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)


def build_presentation(maps: dict, tables: dict) -> Path:
    """
    maps  = {'rgb_compare': Path, 'ndvi_2022': Path, 'lc_compare': Path,
             'ndbi_2022': Path, 'mndwi_2022': Path, 'bsi_2022': Path,
             'dndvi': Path, 'dndbi': Path, 'graph_indices': Path,
             'graph_lc': Path, 'graph_urb_desert': Path, ...}
    tables = {'mean_idx': DataFrame, 'lc_area': DataFrame,
              'change': DataFrame, 'transition': DataFrame}
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    _add_title_slide(
        prs,
        "카타르 월드컵은 사막을 어떻게 바꾸었는가?",
        "2022 FIFA 월드컵 개최에 따른 Lusail 지역 도시화 및 녹지 변화 분석\n"
        "Landsat C2 L2 기반 다시기 원격탐사 분석",
    )

    # 2. Research Need (수업 연계성)
    _add_content_slide(
        prs, "연구 배경 및 필요성 (수업 연계)",
        [
            "• 사막 기후 국가(카타르)의 메가 스포츠 이벤트로 인한 급격한 토지피복 변화 모니터링",
            "• 육안 판독의 주관성을 배제하고 정량적 원격탐사(Remote Sensing) 분석 기법의 필요성",
            "• 센서 및 기하/대기 조건 차이를 극복하기 위한 USGS 대기보정(L2) 데이터 활용",
            "• 스펙트럴 지수 및 비지도 분류(KMeans) 기법을 통한 다중시계열 분석 프로세스 구현",
            "• 사막 도시의 급격한 환경 변화를 입증하여 지구관측 기법의 유효성 검증",
        ],
    )

    # 3. Study Area & CRS
    _add_content_slide(
        prs, "연구지역 및 분석 좌표계",
        [
            f"• 대상 지역: 카타르 루사일 신도시 (중심: {config.CENTER_LAT}°N, {config.CENTER_LON}°E)",
            "• AOI 크기: 약 12 km × 12 km 범위 (루사일 경기장, 마리나, 미개발 사막 포함)",
            "• 분석 투영계: UTM Zone 39N (EPSG:32639) 선정",
            "  - 경위도 좌표계(WGS84)의 위도별 면적 왜곡 한계 극복",
            "  - 평면 직각 투영을 통해 정량적 지표면 면적(km²) 연산 신뢰성 확보",
            "  - 30m 해상도 격자 크기 정합 유지",
        ],
        image_path=maps.get("rgb_2022"),
    )

    # 4. Data Specification
    _add_content_slide(
        prs, "위성 데이터 및 밴드 특성",
        [
            "• 데이터 소스: USGS Landsat Collection 2 Level-2 Surface Reflectance",
            "• 대상 시기: 2000년(사막 Baseline) / 2010년(유치 확정) / 2022년(월드컵 개최)",
            "• 계절 통제: 1~4월(건기) 영상 선별을 통해 일조 조건 및 기후적 변동 오차 최소화",
            "• 사용 밴드: Blue, Green, Red, NIR, SWIR1, SWIR2 대역 물리량 융합",
            "• Landsat 7 SLC-off 결측 복구: 1~4월 내 6개 영상을 대상으로 Median 합성(Compositing)을 적용해 데이터 결측 해소",
        ],
    )

    # 5. Workflow
    _add_content_slide(
        prs, "원격탐사 분석 파이프라인",
        [
            "1. STAC API 검색: 클라우드 커버 < 10% 기준 최적 씬 검색",
            "2. Windowed Read: 대용량 영상 전체가 아닌 12x12 km 영역만 효율적 크롭 리딩",
            "3. 해상도 및 투영 정합: 30 m 격자로 통일하여 스택(Stack) 생성",
            "4. 분광 지수 연산: NDVI, NDBI, MNDWI, BSI 산출",
            "5. 통합 KMeans 분류: 다년도 연합 피처 스페이스 구축 및 분류 실행",
            "6. 전이행렬 및 통계: 연도별 면적 산출 및 2000→2022 변화량 정량화",
            "7. 지도 생산: 방위표, 스케일바, 위경도 좌표 그리드, 범례가 포함된 지도 제작",
        ],
    )

    # 6. Spectral Indices
    _add_content_slide(
        prs, "원격탐사 분광 지수 공식 및 원리",
        [
            "• 식생 지수 (NDVI): (NIR - Red) / (NIR + Red)",
            "  - 식생의 Red 빛 흡수와 NIR 강반사를 활용하여 사막 속 인공 조경 탐지",
            "• 시가지화 지수 (NDBI): (SWIR1 - NIR) / (SWIR1 + NIR)",
            "  - 인공 구조물(콘크리트, 아스팔트)의 SWIR1 강반사 특성을 이용해 도시 면적 강조",
            "• 수정정규수분지수 (MNDWI): (Green - SWIR1) / (Green + SWIR1)",
            "  - 수체의 Green 대역 강반사 및 SWIR1 대역 완전 흡수를 통한 마리나 수역 탐지",
            "• 나지 지수 (BSI): [(SWIR1+Red) - (NIR+Blue)] / [(SWIR1+Red) + (NIR+Blue)]",
            "  - 토양 및 모래사막의 높은 반사도 특성을 모델링하여 사막 면적 변화 추적",
        ],
    )

    # 7. KMeans Classification Method
    _add_content_slide(
        prs, "머신러닝 기반 토지피복 분류",
        [
            "• 10차원 특징 공간 구축: 6개 반사도 밴드 + 4개 분광 지수 결합",
            "  - 단순 반사율 활용 대비 물리적 분류 특성을 극대화하여 지류 구분 성능 향상",
            "• 다년도 연합 KMeans 적합(Fit):",
            "  - 시기 간 클러스터 중심점 및 라벨 불일치를 원천 방지하기 위해 다년도 픽셀을 일괄 샘플링하여 단일 KMeans 모델 피팅 후 개별 시기 예측",
            "• 헝가리안 최적 매칭 알고리즘 적용:",
            "  - 클러스터 중심점 지수 시그니처에 기반한 스코어 행렬을 구축하고 최적 할당(Hungarian Assignment)을 통해 5대 시맨틱 토지피복 라벨을 자동 매핑",
        ],
    )

    # 8. RGB Comparison
    _add_image_only_slide(prs, "True Colour RGB 시계열 비교 (2000 vs 2010 vs 2022)",
                          maps.get("rgb_compare"))

    # 9. NDVI & NDBI Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "변화 탐지 — NDVI(식생) / NDBI(시가지) 차분 분석"
    if maps.get("dndvi") and Path(maps["dndvi"]).exists():
        slide.shapes.add_picture(str(maps["dndvi"]),
                                 Inches(0.3), Inches(1.3), height=Inches(5.6))
    if maps.get("dndbi") and Path(maps["dndbi"]).exists():
        slide.shapes.add_picture(str(maps["dndbi"]),
                                 Inches(6.8), Inches(1.3), height=Inches(5.6))

    # 10. MNDWI & BSI Results
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "수역 및 사막 지수 분석 결과 (MNDWI / BSI 2022)"
    if maps.get("mndwi_2022") and Path(maps["mndwi_2022"]).exists():
        slide.shapes.add_picture(str(maps["mndwi_2022"]),
                                 Inches(0.3), Inches(1.3), height=Inches(5.6))
    if maps.get("bsi_2022") and Path(maps["bsi_2022"]).exists():
        slide.shapes.add_picture(str(maps["bsi_2022"]),
                                 Inches(6.8), Inches(1.3), height=Inches(5.6))

    # 11. Land Cover Comparison
    _add_image_only_slide(prs, "KMeans 토지피복 분류 결과 시계열 비교",
                          maps.get("lc_compare"))

    # 12. Quantitative Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "정량적 면적 변화 및 2000 ➔ 2022 전이행렬"
    df = tables.get("lc_area")
    if df is not None:
        _df_table(slide, df,
                  left=Inches(0.4), top=Inches(1.2),
                  width=Inches(6.2), height=Inches(2.4))
    df2 = tables.get("change")
    if df2 is not None:
        _df_table(slide, df2,
                  left=Inches(0.4), top=Inches(3.9),
                  width=Inches(6.2), height=Inches(2.8))
    if maps.get("graph_lc") and Path(maps["graph_lc"]).exists():
        slide.shapes.add_picture(str(maps["graph_lc"]),
                                 Inches(7.0), Inches(1.2), height=Inches(2.8))
    if maps.get("graph_indices") and Path(maps["graph_indices"]).exists():
        slide.shapes.add_picture(str(maps["graph_indices"]),
                                 Inches(7.0), Inches(4.2), height=Inches(2.8))

    # 13. Conclusions
    _add_content_slide(
        prs, "결론 및 원격탐사적 시사점",
        [
            "✓ 연구 가설 검증: 사막 위 신도시 개발(Urban 증가, Desert 감소) 및 인공 녹지(Vegetation 증가) 실증",
            "✓ 2010~2022 사이 개발 강도가 2000~2010 대비 월등히 높음 (월드컵 유치에 따른 초고속 개발)",
            "✓ 의의: STAC 기반의 효율적인 데이터 파이프라인 및 다차원 원격탐사 융합 분석 프레임워크 구현",
            "✓ 한계점: 비지도 기법(KMeans) 특성에 따른 일부 혼합 표면(Mixed Surface)과 사막 오분류 존재",
            "✓ 향후 개선: Sentinel-2 등 고해상도 자료 결합 및 감독 분류(Random Forest, 딥러닝 U-Net) 도입을 통한 참조 자료 검증",
        ],
    )

    # Speaker notes
    notes = [
        "1. 안녕하십니까, 카타르 월드컵 개최에 따른 루사일 지역의 도시화 및 녹지 변화를 위성 원격탐사로 정량 분석한 기말 프로젝트 발표입니다.",
        "2. 연구 필요성 장표입니다. 단순 육안 판독의 한계를 넘어 대기보정이 완료된 다중시기 위성데이터와 분광 지수, 머신러닝 기법을 활용한 정량 분석 목적을 설정했습니다.",
        "3. 연구지역 및 좌표계 장표입니다. 12x12 km 범위 내 다양한 피복을 담았고, 거리와 면적 왜곡을 없애기 위해 평면 직각 투영인 UTM Zone 39N을 적용했습니다.",
        "4. 위성 데이터 장표입니다. Landsat Collection 2 Level-2 표면반사도를 사용했으며 계절 영향을 통제하기 위해 건기로 고정했고, Landsat 7 결측은 Median Composite로 극복했습니다.",
        "5. 파이프라인 워크플로우입니다. 데이터 API 검색부터 윈도우 크롭, 지수 연산, 통합 KMeans 분류 및 지도 제작까지 전 과정이 자동화된 파이프라인으로 구성되었습니다.",
        "6. 분석 방법 1인 스펙트럴 지수입니다. 식생의 NDVI, 시가지의 NDBI, 수체의 MNDWI, 사막의 BSI 등 타깃 지표 특성을 강조하는 물리적 연산 원리입니다.",
        "7. 분석 방법 2인 KMeans 토지분류입니다. 10차원 특징 공간 구축과 시기 간 일관성을 확보하는 다년도 통합 분류 및 헝가리안 매칭 자동 라벨링을 설명합니다.",
        "8. True Colour RGB 결과입니다. 2000년 완전한 사막 상태에서 2010년 개발 시작, 2022년 월드컵 경기장과 신도시가 완성된 정성적 시계열 비교입니다.",
        "9. 식생 및 시가지 지수 결과입니다. NDVI와 NDBI 지도의 시계열 변화 및 차분 영상을 통해 신도시 중심과 경기장 인근의 강한 도시화 및 녹지화 분포를 확인할 수 있습니다.",
        "10. 수체 및 나지 지수 결과입니다. 신규 수역인 마리나의 출현을 MNDWI로 감지했고, 기존 천연 사막 면적의 대규모 감소를 BSI 감소 패턴으로 실증했습니다.",
        "11. KMeans 토지피복 분류 비교 지도입니다. 2000년 사막 중심 피복에서 2022년 도시 및 혼합 표면, 군데군데 수체와 식생으로 개편된 공간적 변화를 보여줍니다.",
        "12. 정량적 면적 및 전이 행렬 결과입니다. 그래프와 표, 그리고 전이 행렬 수치를 통해 2000년의 사막 면적 중 얼마가 도시와 식생으로 각각 이동했는지 완벽히 입증합니다.",
        "13. 결론 장표입니다. 연구 가설 검증과 분석 파이프라인의 학술적 의의, 그리고 KMeans의 한계점 및 향후 고해상도/감독 분류 융합 등 개선 대책을 제시하며 발표를 마칩니다."
    ]
    for i, slide in enumerate(prs.slides):
        if i < len(notes):
            slide.notes_slide.notes_text_frame.text = notes[i]

    out = config.PPT_DIR / "Lusail_WorldCup_RS_Presentation.pptx"
    prs.save(out)
    return out


def write_qa_document() -> Path:
    out = config.PPT_DIR / "expected_qa.md"
    with open(out, "w", encoding="utf-8") as f:
        f.write("# 예상 질문 & 답변 (10문항)\n\n")
        for q, a in QA_ITEMS:
            f.write(f"## {q}\n\n{a}\n\n")
    return out
